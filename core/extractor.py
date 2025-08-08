"""
Text extraction from PowerPoint presentations.
"""

import logging
import tempfile
import zipfile
from typing import Dict, List, Tuple, Optional, Any
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. Install it for PowerPoint parsing.")

from utils.ocr import OCRProcessor
from utils.helpers import clean_text, extract_numbers_and_dates


class TextExtractor:
    """Extracts text from PowerPoint presentations using multiple methods."""
    
    def __init__(self, ocr_confidence: int = 70):
        """
        Initialize text extractor.
        
        Args:
            ocr_confidence: OCR confidence threshold (0-100)
        """
        self.ocr_processor = OCRProcessor(confidence_threshold=ocr_confidence)
        self.logger = logging.getLogger(__name__)
        
        if not PPTX_AVAILABLE:
            self.logger.error("python-pptx not available. Cannot extract text from PowerPoint files.")
    
    def extract_from_presentation(self, pptx_path: Path) -> Dict[int, Dict[str, Any]]:
        """
        Extract text from all slides in a PowerPoint presentation.
        
        Args:
            pptx_path: Path to the PowerPoint file
            
        Returns:
            Dictionary mapping slide numbers to extracted content
        """
        if not PPTX_AVAILABLE:
            self.logger.error("Cannot extract text: python-pptx not available")
            return {}
        
        try:
            presentation = Presentation(pptx_path)
            slides_data = {}
            
            self.logger.info(f"Extracting text from {len(presentation.slides)} slides")
            
            for slide_index, slide in enumerate(presentation.slides, 1):
                slide_data = self._extract_from_slide(slide, slide_index)
                slides_data[slide_index] = slide_data
                
                self.logger.debug(f"Slide {slide_index}: Extracted {len(slide_data['text'])} text elements")
            
            return slides_data
            
        except Exception as e:
            self.logger.error(f"Failed to extract text from presentation {pptx_path}: {str(e)}")
            return {}
    
    def _extract_from_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """
        Extract text from a single slide.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number (1-based)
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        slide_data = {
            'slide_number': slide_number,
            'text': [],
            'titles': [],
            'body_text': [],
            'table_data': [],
            'shape_text': [],
            'image_text': [],
            'numbers': [],
            'percentages': [],
            'currency': [],
            'dates': [],
            'ocr_used': False,
            'ocr_confidence': 0.0
        }
        
        # Extract text from shapes
        try:
            for shape in slide.shapes:
                shape_text = self._extract_from_shape(shape)
                if shape_text:
                    slide_data['shape_text'].append(shape_text)
                    slide_data['text'].append(shape_text)
        except AttributeError:
            # Some slides may not have shapes attribute
            self.logger.debug(f"Slide {slide_number} has no shapes attribute")
            pass
        
        # Extract text from tables
        try:
            for table in slide.tables:
                table_data = self._extract_from_table(table)
                if table_data:
                    slide_data['table_data'].append(table_data)
                    slide_data['text'].append(table_data)
        except AttributeError:
            # Some slides may not have tables attribute
            self.logger.debug(f"Slide {slide_number} has no tables attribute")
            pass
        
        # Categorize text by type
        self._categorize_text(slide_data)
        
        # Extract structured data
        all_text = ' '.join(slide_data['text'])
        extracted_data = extract_numbers_and_dates(all_text)
        slide_data.update(extracted_data)
        
        # If minimal text found, try OCR on slide images
        if len(slide_data['text']) < 3 and self.ocr_processor.is_available():
            ocr_text = self._extract_with_ocr(slide, slide_number)
            if ocr_text:
                slide_data['image_text'].append(ocr_text)
                slide_data['text'].append(ocr_text)
                slide_data['ocr_used'] = True
        
        return slide_data
    
    def _extract_from_shape(self, shape) -> str:
        """Extract text from a shape."""
        try:
            # Try to get text directly from shape
            if hasattr(shape, 'text') and shape.text:
                return clean_text(shape.text)
            
            # Handle different shape types
            if hasattr(shape, 'shape_type'):
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    if hasattr(shape, 'text_frame'):
                        text_parts = []
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                text_parts.append(paragraph.text.strip())
                        return ' '.join(text_parts)
                
                # Handle other shape types that might contain text
                elif hasattr(shape, 'text_frame'):
                    text_parts = []
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_parts.append(paragraph.text.strip())
                    return ' '.join(text_parts)
            
            return ""
            
        except Exception as e:
            self.logger.debug(f"Failed to extract text from shape: {str(e)}")
            return ""
    
    def _extract_from_table(self, table) -> str:
        """Extract text from a table."""
        try:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text.append(' | '.join(row_text))
            
            return ' | '.join(table_text)
            
        except Exception as e:
            self.logger.debug(f"Failed to extract text from table: {str(e)}")
            return ""
    
    def _categorize_text(self, slide_data: Dict[str, Any]):
        """Categorize extracted text into titles, body text, etc."""
        for text in slide_data['text']:
            if not text:
                continue
            
            # Simple heuristics for categorization
            text_lower = text.lower()
            
            # Check if it's likely a title (short, contains key words)
            if len(text.split()) <= 10 and any(word in text_lower for word in ['title', 'heading', 'slide']):
                slide_data['titles'].append(text)
            else:
                slide_data['body_text'].append(text)
    
    def _extract_with_ocr(self, slide, slide_number: int) -> str:
        """
        Extract text from slide images using OCR.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number for logging
            
        Returns:
            Extracted text from OCR
        """
        try:
            # Create temporary directory for slide images
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                
                # Save slide as image (this would require additional implementation)
                # For now, we'll return empty string as this requires more complex image extraction
                self.logger.debug(f"OCR extraction not fully implemented for slide {slide_number}")
                return ""
                
        except Exception as e:
            self.logger.error(f"OCR extraction failed for slide {slide_number}: {str(e)}")
            return ""
    
    def extract_images_from_pptx(self, pptx_path: Path) -> Dict[int, List[Path]]:
        """
        Extract images from PowerPoint file for OCR processing.
        
        Args:
            pptx_path: Path to the PowerPoint file
            
        Returns:
            Dictionary mapping slide numbers to image paths
        """
        slide_images = {}
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zip_file:
                # PowerPoint files store images in media folder
                media_files = [f for f in zip_file.namelist() if f.startswith('ppt/media/')]
                
                # Group images by slide (this is a simplified approach)
                for media_file in media_files:
                    # Extract image to temporary location
                    with tempfile.TemporaryDirectory() as temp_dir:
                        temp_path = Path(temp_dir)
                        zip_file.extract(media_file, temp_path)
                        image_path = temp_path / media_file
                        
                        # For now, we'll associate all images with slide 1
                        # A more sophisticated approach would be needed to map images to specific slides
                        if 1 not in slide_images:
                            slide_images[1] = []
                        slide_images[1].append(image_path)
            
            return slide_images
            
        except Exception as e:
            self.logger.error(f"Failed to extract images from {pptx_path}: {str(e)}")
            return {}
    
    def get_slide_summary(self, slides_data: Dict[int, Dict[str, Any]]) -> Dict[str, Any]:
        """
        Generate a summary of extracted content.
        
        Args:
            slides_data: Dictionary of slide data
            
        Returns:
            Summary statistics
        """
        total_slides = len(slides_data)
        total_text_elements = sum(len(slide['text']) for slide in slides_data.values())
        slides_with_ocr = sum(1 for slide in slides_data.values() if slide['ocr_used'])
        
        # Collect all numbers, percentages, currency, and dates
        all_numbers = []
        all_percentages = []
        all_currency = []
        all_dates = []
        
        for slide_data in slides_data.values():
            all_numbers.extend(slide_data['numbers'])
            all_percentages.extend(slide_data['percentages'])
            all_currency.extend(slide_data['currency'])
            all_dates.extend(slide_data['dates'])
        
        return {
            'total_slides': total_slides,
            'total_text_elements': total_text_elements,
            'slides_with_ocr': slides_with_ocr,
            'unique_numbers': len(set(all_numbers)),
            'unique_percentages': len(set(all_percentages)),
            'unique_currency': len(set(all_currency)),
            'unique_dates': len(set(all_dates)),
            'avg_text_per_slide': total_text_elements / total_slides if total_slides > 0 else 0
        }
    
    def is_available(self) -> bool:
        """Check if PowerPoint parsing is available."""
        return PPTX_AVAILABLE 